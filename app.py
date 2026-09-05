import io
import re
import zipfile
from pptx import Presentation
import streamlit as st
from PIL import Image, ImageDraw

st.set_page_config(
    page_title="PPT Marked Image Extractor", page_icon="🖼️", layout="centered"
)
st.title("🖼️ PPT Image Extractor (With Drawing Marks)")
st.write(
    "Format: **OutletName_MobileNo_Type_Size.jpg** (Preserves Green/Red Box Marks)"
)

# --- Selection Option ---
image_option = st.radio(
    "Select Image to Export:",
    ("Image 1 (Left / Close View)", "Image 2 (Right / Far View)"),
    index=0
)


def clean_text(text):
    if not text:
        return ""
    clean = re.sub(r'[\\/*?:"<>|\n\r\t]', " ", text)
    clean = re.sub(r"\s+", " ", clean).strip()
    return clean.replace(" ", "_")


def extract_info_from_slide(slide):
    all_text_blocks = []
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip():
            all_text_blocks.append(shape.text_frame.text.strip())
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        all_text_blocks.append(cell.text.strip())

    full_text = "\n".join(all_text_blocks)
    outlet_name, contact_no, media_type, size = "", "", "", ""

    outlet_match = re.search(
        r"Outlet\s*Name\s*[:\-]?\s*([^\n\r]+)", full_text, re.IGNORECASE
    )
    if outlet_match:
        raw_name = outlet_match.group(1).strip()
        cleaned_name = re.split(
            r"Address|City|Contact|Installation|Type|Size|Qty",
            raw_name,
            flags=re.IGNORECASE,
        )[0].strip()
        if cleaned_name:
            outlet_name = cleaned_name

    if not outlet_name:
        ignore_keywords = [
            "qty", "size", "type", "address", "city", "contact",
            "far view", "close view", "board", "installation",
            "dealer_code", "outlet", "before view", "after view"
        ]
        for block in all_text_blocks:
            lines = [l.strip() for l in block.split("\n") if l.strip()]
            for line in lines:
                if not any(k in line.lower() for k in ignore_keywords):
                    if len(line) > 2 and not line.isdigit():
                        outlet_name = line
                        break
            if outlet_name:
                break

    contact_match = re.search(r"\b[6-9]\d{9}\b", full_text)
    if contact_match:
        contact_no = contact_match.group(0)

    type_match = re.search(
        r"\b(NL|FL|BL|SB|GSB|Non-Lit|Flex)\b", full_text, re.IGNORECASE
    )
    if type_match:
        media_type = type_match.group(1).upper()

    size_match = re.search(
        r"(\d{1,3}(?:\.\d+)?\s*x\s*\d{1,3}(?:\.\d+)?)", full_text, re.IGNORECASE
    )
    if size_match:
        size = size_match.group(1).replace(" ", "").lower()

    return outlet_name, contact_no, media_type, size


def get_real_images_from_slide(slide):
    valid_images = []
    for s in slide.shapes:
        if getattr(s, "shape_type", None) == 13 or hasattr(s, "image"):
            try:
                img_bytes = s.image.blob
                test_img = Image.open(io.BytesIO(img_bytes))
                if test_img.width > 50 and test_img.height > 50:
                    valid_images.append(s)
            except Exception:
                continue

    return sorted(valid_images, key=lambda img: img.left)


def process_image_with_marks(slide, img_shape):
    """Overlays ONLY actual inner drawing annotations on the selected image"""
    raw_img_bytes = img_shape.image.blob
    img = Image.open(io.BytesIO(raw_img_bytes)).convert("RGB")

    crop_l = getattr(img_shape, "crop_left", 0) or 0
    crop_r = getattr(img_shape, "crop_right", 0) or 0
    crop_t = getattr(img_shape, "crop_top", 0) or 0
    crop_b = getattr(img_shape, "crop_bottom", 0) or 0

    real_w, real_h = img.size

    if crop_l > 0 or crop_r > 0 or crop_t > 0 or crop_b > 0:
        left_px = int(crop_l * real_w)
        top_px = int(crop_t * real_h)
        right_px = int((1.0 - crop_r) * real_w)
        bottom_px = int((1.0 - crop_b) * real_h)

        if right_px > left_px and bottom_px > top_px:
            img = img.crop((left_px, top_px, right_px, bottom_px))
            real_w, real_h = img.size

    draw = ImageDraw.Draw(img)

    img_left = img_shape.left
    img_top = img_shape.top
    img_width = img_shape.width
    img_height = img_shape.height
    img_right = img_left + img_width
    img_bottom = img_top + img_height

    for s in slide.shapes:
        # Ignore picture shapes, text frames, or non-drawing objects
        if s == img_shape or getattr(s, "shape_type", None) == 13 or s.has_text_frame:
            continue

        s_left = s.left
        s_top = s.top
        s_right = s.left + s.width
        s_bottom = s.top + s.height

        # Strict Filter: Ignore outer slide frames or boxes almost as large as the image itself
        if (s_width := s.width) >= (img_width * 0.85) or (s_height := s.height) >= (img_height * 0.85):
            continue

        # Check if the shape center strictly lies inside the image canvas
        s_center_x = s_left + (s_width / 2)
        s_center_y = s_top + (s_height / 2)

        if img_left <= s_center_x <= img_right and img_top <= s_center_y <= img_bottom:
            rx1 = (s_left - img_left) / img_width
            ry1 = (s_top - img_top) / img_height
            rx2 = (s_right - img_left) / img_width
            ry2 = (s_bottom - img_top) / img_height

            x1 = max(0, rx1 * real_w)
            y1 = max(0, ry1 * real_h)
            x2 = min(real_w, rx2 * real_w)
            y2 = min(real_h, ry2 * real_h)

            mark_color = (0, 255, 0)
            try:
                if hasattr(s, "line") and s.line and s.line.color and s.line.color.rgb:
                    rgb = s.line.color.rgb
                    mark_color = (rgb[0], rgb[1], rgb[2])
            except Exception:
                pass

            stroke = max(4, int(min(real_w, real_h) * 0.012))
            draw.rectangle([x1, y1, x2, y2], outline=mark_color, width=stroke)

    out = io.BytesIO()
    img.save(out, format="JPEG", quality=95)
    return out.getvalue()


uploaded_file = st.file_uploader("Upload PowerPoint File (.pptx)", type=["pptx"])

if uploaded_file is not None:
    st.info(f"📁 **Uploaded File:** {uploaded_file.name} ({round(uploaded_file.size / (1024 * 1024), 2)} MB)")
    
    # --- START BUTTON ---
    if st.button("▶️ Start Extraction", type="primary", use_container_width=True):
        prs = Presentation(uploaded_file)
        zip_buffer = io.BytesIO()

        with st.spinner("Extracting image markings cleanly..."):
            with zipfile.ZipFile(zip_buffer, "a", zipfile.ZIP_DEFLATED, False) as zip_file:
                for i, slide in enumerate(prs.slides):
                    outlet_name, contact_no, media_type, size = extract_info_from_slide(slide)
                    
                    valid_images = get_real_images_from_slide(slide)

                    if valid_images:
                        target_pic = None
                        
                        if "Image 1" in image_option:
                            target_pic = valid_images[0]
                        elif "Image 2" in image_option:
                            if len(valid_images) >= 2:
                                target_pic = valid_images[1]
                            elif len(valid_images) == 1:
                                target_pic = valid_images[0]

                        if target_pic:
                            final_bytes = process_image_with_marks(slide, target_pic)

                            if not outlet_name:
                                outlet_name = f"Slide_{i+1}"

                            components = [clean_text(outlet_name)]
                            if contact_no:
                                components.append(clean_text(contact_no))
                            if media_type:
                                components.append(clean_text(media_type))
                            if size:
                                components.append(clean_text(size))

                            final_name = f"{'_'.join(components)}.jpg"
                            zip_file.writestr(final_name, final_bytes)

        st.success("🎉 Process Complete!")
        st.download_button(
            label="📥 Download Selected Images (ZIP)",
            data=zip_buffer.getvalue(),
            file_name="Renamed_Images.zip",
            mime="application/zip",
        )
